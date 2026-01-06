# =========================================================
# SMART PDF / PPT CHATBOT (OFFLINE, NO LANGCHAIN)
# =========================================================

from PyPDF2 import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
import json
import os
import re

# ================= CONFIG =================

FILE_PATH = r"D:\Desktop\FULL Stack\FSDL TY COMP MERN Stack.pptx"

TEXT_FILE = "document_text.txt"
JSON_FILE = "embeddings.json"

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
CONFIDENCE_SCORE = 0.45
MAX_SENTENCES = 10    # controls answer size

model = SentenceTransformer("all-MiniLM-L6-v2")

used_chunks = set()
question_count = 0

# =========================================================
# TEXT SPLITTER
# =========================================================
def split_text(text, chunk_size, overlap):
    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end].strip())
        start = end - overlap

    return chunks

# =========================================================
# EXTRACT TEXT
# =========================================================
def extract_text():
    text = ""

    if FILE_PATH.lower().endswith(".pdf"):
        reader = PdfReader(FILE_PATH)
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"

    elif FILE_PATH.lower().endswith(".pptx"):
        prs = Presentation(FILE_PATH)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    else:
        raise ValueError("Unsupported file format")

    with open(TEXT_FILE, "w", encoding="utf-8") as f:
        f.write(text)

# =========================================================
# CREATE EMBEDDINGS
# =========================================================
def create_embeddings():
    with open(TEXT_FILE, "r", encoding="utf-8") as f:
        text = f.read()

    chunks = split_text(text, CHUNK_SIZE, CHUNK_OVERLAP)
    embeddings = model.encode(chunks)
    embeddings = embeddings / np.linalg.norm(embeddings, axis=1, keepdims=True)

    with open(JSON_FILE, "w", encoding="utf-8") as f:
        json.dump({"chunks": chunks, "embeddings": embeddings.tolist()}, f, indent=4)

# =========================================================
# LOAD EMBEDDINGS
# =========================================================
def load_embeddings():
    with open(JSON_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)

    return data["chunks"], np.array(data["embeddings"])

# =========================================================
# CLEAN + SHORT ANSWER GENERATOR
# =========================================================
def summarize_text(text):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 40]
    return "\n• ".join(sentences[:MAX_SENTENCES])

# =========================================================
# ANSWER QUESTION (NON-REPETITIVE)
# =========================================================
def answer_question(question, chunks, embeddings):
    q_embedding = model.encode(question)
    q_embedding = q_embedding / np.linalg.norm(q_embedding)

    scores = np.dot(embeddings, q_embedding)

    ranked = sorted(
        zip(chunks, scores),
        key=lambda x: x[1],
        reverse=True
    )

    for chunk, score in ranked:
        if score < CONFIDENCE_SCORE:
            break

        if chunk not in used_chunks:
            used_chunks.add(chunk)
            return summarize_text(chunk)

    return "No new relevant information found."

# =========================================================
# MAIN
# =========================================================
if __name__ == "__main__":

    if not os.path.exists(JSON_FILE):
        print("🔹 Extracting document text...")
        extract_text()

        print("🔹 Creating embeddings...")
        create_embeddings()

    chunks, embeddings = load_embeddings()

    print("\n📘 SMART DOCUMENT CHATBOT READY")
    print("Type 'exit' to quit\n")

    while True:
        question = input(">>> Ask Question: ")

        if question.lower() == "exit":
            break

        question_count += 1
        print("\n" + "=" * 60)
        print(f"Q{question_count}: {question}")
        print("-" * 60)

        answer = answer_question(question, chunks, embeddings)
        print("Answer:")
        print("• " + answer)
        print("=" * 60 + "\n")
